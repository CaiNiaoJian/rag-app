"""text_coverage 去假回归（TODO「完整性指标目前是假的」项）。

历史问题：docx/pptx/xlsx 走 else 分支恒写 1.0——docx 脚注静默丢字、SmartArt 全丢时
仪表盘照样满分。修复后的契约：

- 分母来自**不经过解析器**的旁路（docx/pptx 读包内 XML；xlsx 扫描时累计；PDF 原有）；
- 无损解析 → 恰好 1.0；有丢失 → 按比例下降；分母统计不了 → None（UI「—」），
  决不能再回退成假 1.0。
"""

from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from docfactory.config import Settings
from docfactory.parsers import parse_document
from docfactory.parsers.raw_text import raw_char_total

# ---------------------------------------------------------------- 端到端


def _parse(src: Path, paths, fmt: str):
    return parse_document(
        src=src, doc_id=f"cov-{src.stem}", fmt=fmt, paths=paths, settings=Settings(), ctx=None
    )


def test_docx_lossless_parse_scores_exactly_one(paths, tmp_path: Path) -> None:
    """标题/列表/表格/页眉页脚全被解析器抽走 → 覆盖率应恰好 1.0（真算而非假 1.0）。"""
    from docx import Document

    doc = Document()
    doc.sections[0].header.paragraphs[0].text = "内部资料 · 请勿外传"
    doc.sections[0].footer.paragraphs[0].text = "第 1 页"
    doc.add_heading("第1章 总则", level=1)
    doc.add_paragraph("本合同由甲乙双方在平等自愿基础上签订。")
    doc.add_paragraph("源代码及构建脚本", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "项目"
    table.cell(0, 1).text = "金额"
    table.cell(1, 0).text = "服务器"
    table.cell(1, 1).text = "120000"
    src = tmp_path / "lossless.docx"
    doc.save(str(src))

    ir = _parse(src, paths, "docx")
    assert ir.doc.metrics.text_coverage == 1.0


def test_docx_footnote_loss_is_visible(paths, tmp_path: Path) -> None:
    """脚注是当前解析器的已知盲区（TODO「docx 脚注静默丢字」）：
    指标去假后，这份丢失必须体现为覆盖率 < 1，而不是继续满分。"""
    from docx import Document

    doc = Document()
    body_text = "合同正文内容共计二十个可见字符整"
    doc.add_paragraph(body_text)
    src = tmp_path / "with_footnote.docx"
    doc.save(str(src))

    # python-docx 不支持写脚注：直接向包内追加 footnotes 部件。
    # 解析器不读它（这正是被测的丢失路径），raw_text 旁路统计会读到。
    footnote_text = "关键限定条件往往写在脚注里一旦丢失后果严重"
    footnotes_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:footnote w:id="1"><w:p><w:r><w:t>' + footnote_text + "</w:t></w:r></w:p></w:footnote>"
        "</w:footnotes>"
    )
    with zipfile.ZipFile(src, "a") as zf:
        zf.writestr("word/footnotes.xml", footnotes_xml)

    ir = _parse(src, paths, "docx")
    cov = ir.doc.metrics.text_coverage
    assert cov is not None and cov < 1.0, "脚注文本丢失时覆盖率不得再显示满分"
    expected = len(body_text) / (len(body_text) + len(footnote_text))
    assert abs(cov - expected) < 0.02


def test_pptx_lossless_parse_scores_exactly_one(paths, tmp_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "季度经营回顾"
    slide.placeholders[1].text = "营收同比增长 23%"
    slide.notes_slide.notes_text_frame.text = "强调增长主要来自政企客户。"
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "区域分布"
    tbl = slide2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    tbl.cell(0, 0).text = "区域"
    tbl.cell(0, 1).text = "占比"
    tbl.cell(1, 0).text = "华东"
    tbl.cell(1, 1).text = "45%"
    src = tmp_path / "deck.pptx"
    prs.save(str(src))

    ir = _parse(src, paths, "pptx")
    assert ir.doc.metrics.text_coverage == 1.0


def test_xlsx_coverage_accumulated_during_scan(paths, tmp_path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "预算"
    ws.append(["科目", "预算"])
    ws.append(["人力", 1200000])
    ws["A4"] = "=SUM(B2:B3)"  # 无缓存值的公式：回退原文，不允许把覆盖率打穿
    src = tmp_path / "book.xlsx"
    wb.save(str(src))

    ir = _parse(src, paths, "xlsx")
    assert ir.doc.metrics.text_coverage == 1.0


def test_coverage_is_none_when_denominator_unavailable(
    paths, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """分母统计失败 → 诚实记 None；回退成 1.0 就是回到造假老路。"""
    from docx import Document

    import docfactory.parsers.raw_text as raw_text

    doc = Document()
    doc.add_paragraph("一段正文")
    src = tmp_path / "plain.docx"
    doc.save(str(src))

    monkeypatch.setattr(raw_text, "raw_char_total", lambda *_: None)
    ir = _parse(src, paths, "docx")
    assert ir.doc.metrics.text_coverage is None


# ---------------------------------------------------------------- raw_text 单元


def _docx_zip(path: Path, entries: dict[str, str]) -> Path:
    with zipfile.ZipFile(path, "w") as zf:
        for name, xml in entries.items():
            zf.writestr(name, xml)
    return path


def _wp(text: str) -> str:
    return f"<w:p><w:r><w:t>{text}</w:t></w:r></w:p>"


_W_OPEN = '<w:{root} xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'


def test_raw_docx_counts_footnotes_and_dedups_headers(tmp_path: Path) -> None:
    doc = _W_OPEN.format(root="document") + "<w:body>" + _wp("正文十个字符长度計") + "</w:body></w:document>"
    foot = _W_OPEN.format(root="footnotes") + '<w:footnote w:id="1">' + _wp("脚注四字") + "</w:footnote></w:footnotes>"
    hdr = _W_OPEN.format(root="hdr") + _wp("同一份页眉") + "</w:hdr>"
    ftr = _W_OPEN.format(root="ftr") + _wp("页脚三") + "</w:ftr>"
    src = _docx_zip(tmp_path / "t.docx", {
        "word/document.xml": doc,
        "word/footnotes.xml": foot,
        "word/header1.xml": hdr,
        "word/header2.xml": hdr,   # 与 header1 完全相同：解析器同文去重，分母也只计一次
        "word/footer1.xml": ftr,
    })
    # 9（正文）+ 4（脚注）+ 5（页眉去重后）+ 3（页脚）
    assert raw_char_total(src, "docx") == 9 + 4 + 5 + 3


_A_NS = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
_P_NS = 'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'


def test_raw_pptx_excludes_fld_and_counts_notes_body_and_diagram(tmp_path: Path) -> None:
    slide = (
        f"<p:sld {_P_NS} {_A_NS}><p:cSld><p:spTree>"
        "<p:sp><p:txBody><a:p><a:r><a:t>标题四个字</a:t></a:r>"
        '<a:fld id="{X}" type="slidenum"><a:t>3</a:t></a:fld>'  # 页码域：不计
        "</a:p></p:txBody></p:sp>"
        "</p:spTree></p:cSld></p:sld>"
    )
    notes = (
        f"<p:notes {_P_NS} {_A_NS}><p:cSld><p:spTree>"
        '<p:sp><p:nvSpPr><p:nvPr><p:ph type="body"/></p:nvPr></p:nvSpPr>'
        "<p:txBody><a:p><a:r><a:t>备注正文六字</a:t></a:r></a:p></p:txBody></p:sp>"
        '<p:sp><p:nvSpPr><p:nvPr><p:ph type="sldNum"/></p:nvPr></p:nvSpPr>'
        "<p:txBody><a:p><a:r><a:t>1</a:t></a:r></a:p></p:txBody></p:sp>"  # 页码占位：不计
        "</p:spTree></p:cSld></p:notes>"
    )
    diagram = f"<dgm:dataModel xmlns:dgm=\"http://schemas.openxmlformats.org/drawingml/2006/diagram\" {_A_NS}><a:t>图形文本</a:t></dgm:dataModel>"
    src = _docx_zip(tmp_path / "t.pptx", {
        "ppt/slides/slide1.xml": slide,
        "ppt/notesSlides/notesSlide1.xml": notes,
        "ppt/diagrams/data1.xml": diagram,
    })
    # 5（正文，页码域剔除）+ 6（备注 body）+ 4（SmartArt）
    assert raw_char_total(src, "pptx") == 5 + 6 + 4


def test_raw_text_unsupported_format_returns_none(tmp_path: Path) -> None:
    src = tmp_path / "x.bin"
    src.write_bytes(b"not a zip")
    assert raw_char_total(src, "docx") is None
    assert raw_char_total(src, "pdf") is None
