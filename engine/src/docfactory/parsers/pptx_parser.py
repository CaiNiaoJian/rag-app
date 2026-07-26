"""pptx 解析器（python-pptx，MIT；03 章 §1 直解路径）。

每张 slide 一个 ``slide`` 节点（title/notes 进 content），页内形状按**版面位置**排序后
作为子节点：文本框 → paragraph、表格 → table、图片 → figure（落 assets）、组合 → 递归展开。

两个设计取舍：

- **按位置排序而不是按 XML 顺序**：pptx 里形状的 XML 顺序是「创建/层叠顺序」，
  跟人阅读的顺序无关（后加的标题可能排在最后）。按 (top, left) 排序更接近阅读序，
  也让切片层「每 slide 一片」的文本读起来是顺的（04 章 §3.1）。
- **slide 序号即 prov.page**：PPT 天然一页一 slide，页码没有歧义（与 docx 相反）。

降级：结构化解析异常时上层调 ``parse_text_fallback()``，退成逐页纯文本（L2）。
"""

from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docfactory.errors import DocFactoryError
from docfactory.ir import IRNode, NodeContent, Prov, TableCell, TableContent
from docfactory.parsers import ParseEnv, assert_ooxml_readable

# 排序兜底值：形状缺 top/left（继承版式占位符）时排到最后，避免 None 参与比较
_FAR = 10**9


def _open(src: Path) -> Any:
    assert_ooxml_readable(src)
    try:
        return Presentation(str(src))
    except DocFactoryError:
        raise
    except Exception as exc:
        raise DocFactoryError("E01", f"打开 pptx 失败：{src.name}（{type(exc).__name__}: {exc}）") from exc


def _pos(shape: Any) -> tuple[int, int]:
    try:
        top = shape.top if shape.top is not None else _FAR
        left = shape.left if shape.left is not None else _FAR
    except (AttributeError, ValueError):
        return _FAR, _FAR
    return int(top), int(left)


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """按阅读序展开形状树：组合形状递归进去，组合内部同样按位置排序。"""
    for shape in sorted(shapes, key=_pos):
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _notes(slide: Any) -> str:
    """演讲者备注（往往是全篇信息密度最高的文字，务必并入）。"""
    try:
        if not slide.has_notes_slide:
            return ""
        frame = slide.notes_slide.notes_text_frame
        return (frame.text or "").strip() if frame is not None else ""
    except Exception:
        return ""


def _title(slide: Any) -> tuple[str, Any]:
    """返回 (标题文本, 标题占位符形状)；后者用于避免标题在正文里重复出现。"""
    try:
        placeholder = slide.shapes.title
    except Exception:
        return "", None
    if placeholder is None:
        return "", None
    try:
        return (placeholder.text or "").strip(), placeholder
    except Exception:
        return "", placeholder


def _table_content(table: Any) -> TableContent:
    """pptx 表格 → IR cells。

    python-pptx 的合并单元格模型比 docx 直白：``is_merge_origin`` 的格带 span_height/width，
    被覆盖的格 ``is_spanned=True`` 直接跳过（IR 只保留合并区左上格）。
    第一行是否表头按 ``first_row`` 样式开关判定，取不到则默认首行为表头。
    """
    try:
        header_first = bool(table.first_row)
    except Exception:
        header_first = True
    cells: list[TableCell] = []
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            if getattr(cell, "is_spanned", False):
                continue
            rowspan = int(getattr(cell, "span_height", 1) or 1)
            colspan = int(getattr(cell, "span_width", 1) or 1)
            cells.append(TableCell(
                r=r, c=c, rowspan=rowspan, colspan=colspan,
                text=(cell.text or "").strip(),
                is_header=header_first and r == 0,
            ))
    return TableContent(cells=cells)


def _image(shape: Any) -> tuple[bytes, str] | None:
    try:
        image = shape.image
        return image.blob, f".{(image.ext or 'png').lstrip('.')}"
    except Exception:  # 链接式图片（外部文件）没有 blob，跳过
        return None


def _shape_paragraphs(shape: Any) -> list[str]:
    """文本框里逐段取文本：一段一个 paragraph 节点，保留作者的分点结构。"""
    out: list[str] = []
    try:
        frame = shape.text_frame
    except Exception:
        return out
    for para in frame.paragraphs:
        text = "".join(run.text or "" for run in para.runs).strip()
        if not text:  # runs 为空但有文本（域/自动编号）时退回整段文本
            text = (para.text or "").strip()
        if text:
            out.append(text)
    return out


def _emit_shape(shape: Any, slide_node: IRNode, page: int, env: ParseEnv) -> None:
    prov = [Prov(page=page)]
    builder = env.builder

    if getattr(shape, "has_table", False):
        table = _table_content(shape.table)
        if table.cells:
            # IRNode.content 恒为 NodeContent；TableContent 要包一层，直接传会被 pydantic 拒收
            builder.add("table", parent=slide_node,
                        content=NodeContent(table=table), prov=prov)
        return

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        got = _image(shape)
        if got is not None:
            blob, ext = got
            ref = env.save_asset(blob, ext)
            builder.add(
                "figure", parent=slide_node,
                content=NodeContent(image_ref=ref, ocr=False, caption=_shape_name(shape)),
                prov=prov,
            )
        return

    if getattr(shape, "has_chart", False):
        # 图表 V1 只留标题（03 章 §7 同款取舍：不还原图形，避免产出误导性数据）
        title = _chart_title(shape)
        if title:
            builder.add(
                "paragraph", parent=slide_node,
                content=NodeContent(text=f"[图表] {title}"), prov=prov,
            )
        return

    if getattr(shape, "has_text_frame", False):
        for text in _shape_paragraphs(shape):
            builder.add("paragraph", parent=slide_node, content=NodeContent(text=text), prov=prov)


def _shape_name(shape: Any) -> str | None:
    """图片的替代文字（alt text）作 caption。

    只取 alt text，**不退回 shape.name** —— 后者默认是 "Picture 3" 这类自动名，
    当图注写进 Markdown/数据集只会是噪声；宁可没有 caption。
    """
    try:
        alt = (shape._element._nvXxPr.cNvPr.get("descr") or "").strip()
    except Exception:
        return None
    return alt or None


def _chart_title(shape: Any) -> str:
    try:
        chart = shape.chart
        if chart.has_title and chart.chart_title.has_text_frame:
            return (chart.chart_title.text_frame.text or "").strip()
    except Exception:
        return ""
    return ""


def _same_shape(a: Any, b: Any) -> bool:
    """判断两个形状是否是同一个。

    python-pptx 每次访问都返回**新的代理对象**（``slide.shapes.title`` 与遍历得到的
    同一形状不是同一个 Python 对象），用 ``is`` 比较永远为 False —— 标题会被重复
    写成一个 paragraph。所以比底层 XML 元素的身份。
    """
    if b is None:
        return False
    ea, eb = getattr(a, "_element", a), getattr(b, "_element", b)
    return ea is eb


def parse(src: Path, env: ParseEnv) -> None:
    prs = _open(src)
    slides = list(prs.slides)
    total = max(1, len(slides))
    env.page_count = len(slides)

    for page, slide in enumerate(slides, start=1):
        title, title_shape = _title(slide)
        slide_node = env.builder.add(
            "slide",
            content=NodeContent(title=title or None, notes=_notes(slide) or None),
            prov=[Prov(page=page)],
        )
        for shape in _iter_shapes(slide.shapes):
            if _same_shape(shape, title_shape):
                continue  # 标题已在 slide.content.title 里，不重复成段
            _emit_shape(shape, slide_node, page, env)

        env.mark_level(page, "L0")
        env.progress(page, total)
        env.check_cancel()


def parse_text_fallback(src: Path, env: ParseEnv) -> None:
    """L2 兜底：每页一个 slide 节点 + 一段合并文本，放弃形状结构。"""
    prs = _open(src)
    slides = list(prs.slides)
    env.page_count = len(slides)
    for page, slide in enumerate(slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = (shape.text_frame.text or "").strip()
                if text:
                    texts.append(text)
        node = env.builder.add(
            "slide", content=NodeContent(notes=_notes(slide) or None), prov=[Prov(page=page)]
        )
        if texts:
            env.builder.add(
                "paragraph", parent=node,
                content=NodeContent(text="\n".join(texts)), prov=[Prov(page=page)],
            )
        env.mark_level(page, "L2")
