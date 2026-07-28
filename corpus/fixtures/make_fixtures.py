"""程序化生成 golden corpus 的基础样本（08 章 §3.1）。

为什么不用真实文档：真实语料有版权与体积问题，不能入库；但 CI 必须有可跑的最小回归集。
这里用 python-docx / python-pptx / openpyxl / reportlab 造出**结构确定**的样本 ——
每个样本刻意包含一个已知的失败模式（合并单元格、跨级标题、公式缓存值、超大表截断……），
标注就是生成脚本的意图，永远不会与样本漂移。

真实语料放 corpus/samples/（.gitignore 忽略），与这里的产物一起被 run_corpus.py 扫到。

用法：python corpus/fixtures/make_fixtures.py [输出目录]
"""

from __future__ import annotations

import struct
import sys
import zipfile
import zlib
from collections.abc import Callable
from pathlib import Path

DEFAULT_OUT = Path(__file__).resolve().parent.parent / "samples"


# ---------------------------------------------------------------- 公共工具


def _register_cjk_font() -> str:
    """注册一枚中文 TTF，返回可用字体名；机器上没有中文字体时返回 "Helvetica"。

    调用方据此决定写中文还是英文正文 —— 目的是让这个脚本在任何机器上都能跑完，
    而不是在缺字体的机器上生成一堆问号。
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    if "CJK" in pdfmetrics.getRegisteredFontNames():
        return "CJK"
    for candidate in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\simhei.ttf"):
        if Path(candidate).exists():
            try:
                pdfmetrics.registerFont(TTFont("CJK", candidate))
                return "CJK"
            except Exception:
                continue
    return "Helvetica"


def _png_bytes(width: int, height: int, rgb: tuple[int, int, int]) -> bytes:
    """手写一张纯色 PNG。

    刻意不引 Pillow：它不是引擎的运行时依赖，为了造 fixture 而把它拉进 dev 组，
    等于让 corpus 的可跑性依赖一个与产品无关的包。PNG 的最小结构就三个 chunk，
    直接写比解释为什么多一个依赖更省事。
    """
    raw = b"".join(b"\x00" + bytes(rgb) * width for _ in range(height))

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    return (
        b"\x89PNG\r\n\x1a\n"
        # 位深 8、颜色类型 2（truecolor RGB），无隔行
        + chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


# ---------------------------------------------------------------- docx


def make_headings_docx(path: Path) -> None:
    """多级标题 + 列表：验证层级还原与 list 节点。"""
    from docx import Document

    doc = Document()
    doc.add_heading("第1章 总则", level=1)
    doc.add_paragraph("本合同由甲乙双方于平等自愿基础上签订，具备完全法律效力。")
    doc.add_heading("1.1 定义", level=2)
    doc.add_paragraph("本合同中「交付物」指乙方按约定向甲方提供的全部成果。")
    doc.add_heading("1.1.1 交付物范围", level=3)
    doc.add_paragraph("交付物包括但不限于：", style=None)
    for item in ("源代码及构建脚本", "部署文档与运维手册", "验收测试报告"):
        doc.add_paragraph(item, style="List Bullet")
    doc.add_heading("第2章 交付条款", level=1)
    doc.add_paragraph("交付期限为合同签订后 90 天。逾期按日千分之三支付违约金。")
    doc.save(path)


def make_merged_table_docx(path: Path) -> None:
    """含合并单元格的表格：验证 rowspan/colspan 与取值归左上格。"""
    from docx import Document

    doc = Document()
    doc.add_heading("采购清单", level=1)
    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ("项目", "规格", "金额")
    for i, text in enumerate(headers):
        table.cell(0, i).text = text
    rows = [
        ("服务器", "2U 机架式", "120000"),
        ("交换机", "48 口千兆", "18000"),
        ("合计", "", "138000"),
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            table.cell(r, c).text = text
    # 「合计」行横跨前两列：这是解析器最容易漏掉的一类结构
    table.cell(3, 0).merge(table.cell(3, 1))
    doc.add_paragraph("以上金额均为含税价。")
    doc.save(path)


def make_empty_docx(path: Path) -> None:
    """几乎空的文档：验证 E05「未能提取到有效内容」的优雅报错路径。"""
    from docx import Document

    Document().save(path)


def make_bilingual_docx(path: Path) -> None:
    """中英混排 + 非 ASCII 符号：验证编码全链路（IR→MD→CSV）不出现替换字符。

    覆盖 08 章 §3.1 的「中英混排」场景。特意混进全角标点、上下标、℃、①、
    以及日文假名——这类字符在「Windows 控制台 GBK / CSV 无 BOM / MD 转义」
    三处都掉过链子，放进回归集才不会靠人记得去试。
    """
    from docx import Document

    doc = Document()
    doc.add_heading("Bilingual Test 中英混排测试", level=1)
    doc.add_paragraph(
        "The system supports 七种格式 including PDF、DOCX 和 XLSX，"
        "并在 offline 环境下运行（no external services）。"
    )
    doc.add_heading("2. 技术参数 Technical Specs", level=2)
    doc.add_paragraph("工作温度：-20℃ ~ +60℃，湿度 ≤ 85%（无凝露）。")
    doc.add_paragraph("面积单位 m²，体积单位 m³，误差 ±0.5%。")
    doc.add_paragraph("① 第一项 First item　② 第二项 Second item")
    doc.add_paragraph("日本語のテキストも含まれています。Mixed with English.")
    doc.add_paragraph("引号测试：「全角括号」『日式引号』“English quotes”")
    doc.save(path)


def make_nested_list_docx(path: Path) -> None:
    """三级嵌套列表：验证 list 节点的层级还原（04 章 §2 list/list_item）。

    扁平化处理列表是很常见的偷懒做法——三级结构一旦被压平，切片时
    「第 3 条下面的两个子项」就会跟第 4 条粘在一起，语义直接错位。
    """
    from docx import Document

    doc = Document()
    doc.add_heading("实施步骤", level=1)
    doc.add_paragraph("环境准备", style="List Number")
    doc.add_paragraph("检查操作系统版本", style="List Bullet 2")
    doc.add_paragraph("确认磁盘剩余空间不少于 10GB", style="List Bullet 2")
    doc.add_paragraph("确认 CPU 支持 AVX2 指令集", style="List Bullet 3")
    doc.add_paragraph("安装部署", style="List Number")
    doc.add_paragraph("运行安装程序", style="List Bullet 2")
    doc.add_paragraph("导入历史文档", style="List Bullet 2")
    doc.add_paragraph("验收测试", style="List Number")
    doc.add_paragraph("以上步骤须按顺序执行。")
    doc.save(path)


def make_header_footer_docx(path: Path) -> None:
    """页眉页脚：验证 header/footer 节点识别与 drop_header_footer 处理选项。

    页眉页脚在每一页重复，不剔除就会在切片里反复出现同一串「公司机密」，
    既污染 RAG 语料又浪费 token —— FR-16 把它做成开关正是为此。
    """
    from docx import Document

    doc = Document()
    section = doc.sections[0]
    section.header.paragraphs[0].text = "内部资料 · 请勿外传"
    section.footer.paragraphs[0].text = "DocFactory 技术文档 —— 第 1 页"

    doc.add_heading("系统概述", level=1)
    doc.add_paragraph("本系统用于内网环境下的文档批量解析与数据集构建。")
    doc.add_heading("部署要求", level=2)
    doc.add_paragraph("需要 Windows 10 1809 及以上版本，内存不少于 8GB。")
    doc.save(path)


def make_image_docx(path: Path) -> None:
    """内嵌图片：验证图片抽取落 assets 目录与 MD 的相对路径引用（05 章 §1）。"""
    from io import BytesIO

    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_heading("架构图示", level=1)
    doc.add_paragraph("下图为系统整体架构：")
    doc.add_picture(BytesIO(_png_bytes(240, 160, (60, 90, 200))), width=Inches(3))
    doc.add_paragraph("图 1 系统架构图")
    doc.add_paragraph("图中蓝色部分为本次交付范围。")
    doc.add_picture(BytesIO(_png_bytes(160, 160, (200, 120, 40))), width=Inches(2))
    doc.add_paragraph("图 2 部署拓扑")
    doc.save(path)


# ---------------------------------------------------------------- pptx


def make_slides_pptx(path: Path) -> None:
    """图文混排 PPT + 演讲者备注：验证 slide 节点与 notes 并入。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "季度经营回顾"
    slide.placeholders[1].text = "营收同比增长 23%\n毛利率提升 4 个百分点"
    slide.notes_slide.notes_text_frame.text = "此处强调增长主要来自政企客户。"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "区域分布"
    table = slide2.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "区域"
    table.cell(0, 1).text = "占比"
    table.cell(1, 0).text = "华东"
    table.cell(1, 1).text = "45%"
    table.cell(2, 0).text = "华南"
    table.cell(2, 1).text = "31%"
    prs.save(path)


def make_large_deck_pptx(path: Path) -> None:
    """210 页大 PPT：08 章 §3.1 的「200+ 页大 PPT」场景。

    这份样本量的是**规模**不是结构：解析器在几页上正确不代表在两百页上不退化
    （每页新建 IR 节点、切片按页聚合、进度按页上报都在这里才暴露线性度问题）。
    每页内容刻意可预测，标注才能断言「第 N 页确实在」而不是只数个总数。
    """
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]  # 标题 + 内容
    for i in range(1, 211):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"第 {i} 节 业务单元 {i:03d}"
        slide.placeholders[1].text = (
            f"本节介绍业务单元 {i:03d} 的运营情况。\n"
            f"本期营收 {i * 137 % 9000 + 1000} 万元。\n"
            f"负责人：管理员{i % 17}"
        )
    prs.save(path)


# ---------------------------------------------------------------- xlsx


def make_multisheet_xlsx(path: Path) -> None:
    """多 sheet + 合并单元格 + 公式缓存值 + 空 sheet：Excel 契约（03 章 §7）的主要分支。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "预算"
    ws["A1"] = "部门预算表"
    ws.merge_cells("A1:C1")
    for col, text in zip("ABC", ("科目", "预算", "实际"), strict=True):
        ws[f"{col}2"] = text
    data = [("人力", 1200000, 1180000), ("市场", 500000, 620000), ("研发", 2000000, 1950000)]
    for i, (name, budget, actual) in enumerate(data, start=3):
        ws[f"A{i}"], ws[f"B{i}"], ws[f"C{i}"] = name, budget, actual
    # 公式单元格：openpyxl 写入时没有缓存值，解析器应回退公式原文并记 warning
    ws["B6"] = "=SUM(B3:B5)"
    ws["A6"] = "合计"

    detail = wb.create_sheet("明细")
    detail["A1"] = "日期"
    detail["B1"] = "金额"
    for i in range(2, 12):
        detail[f"A{i}"] = f"2026-07-{i - 1:02d}"
        detail[f"B{i}"] = i * 1000

    wb.create_sheet("空表")  # 空 sheet 应被跳过
    wb.save(path)


def make_wide_xlsx(path: Path) -> None:
    """较多行的表：验证切片层按行组切分并复制表头（04 章 §3.1 规则 1）。"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "流水"
    ws.append(["序号", "客户", "品类", "数量", "单价", "金额"])
    for i in range(1, 601):
        ws.append([i, f"客户{i % 37}", ["硬件", "软件", "服务"][i % 3], i % 9 + 1,
                   round(100 + i * 0.7, 2), round((i % 9 + 1) * (100 + i * 0.7), 2)])
    wb.save(path)


def make_legacy_table_xls(path: Path) -> None:
    """BIFF 旧格式 .xls：验证 xlrd 兜底转换路径（03 章 §4 备选路径）。

    CI 环境没有 LibreOffice，这份样本走的正是 xlrd → xlsx 的纯 Python 兜底——
    合并标题、加粗表头、日期列都是转换器必须穿透保留的信息。
    本机装有 LibreOffice 时则走 soffice 主路径，两条路径的断言口径相同。
    """
    import xlwt

    wb = xlwt.Workbook()
    ws = wb.add_sheet("预算")
    bold = xlwt.easyxf("font: bold on")
    ymd = xlwt.easyxf(num_format_str="yyyy-mm-dd")

    ws.write_merge(0, 0, 0, 2, "部门预算表")
    for col, text in enumerate(("科目", "预算", "签约日")):
        ws.write(1, col, text, bold)
    from datetime import date as _date
    rows = [
        ("人力", 1200000, _date(2026, 2, 11)),
        ("市场", 500000, _date(2026, 3, 5)),
        ("研发", 2000000, _date(2026, 3, 28)),
    ]
    for r, (name, budget, day) in enumerate(rows, start=2):
        ws.write(r, 0, name)
        ws.write(r, 1, budget)
        ws.write(r, 2, day, ymd)

    detail = wb.add_sheet("明细")
    detail.write(0, 0, "备注")
    detail.write(1, 0, "以上金额均为含税价")
    wb.save(str(path))


def make_sparse_regions_xlsx(path: Path) -> None:
    """一个 sheet 内多个不连续数据区域：验证 sheet_region 切分（03 章 §7）。

    Excel 的实际用法很少是「一个 sheet 一张干净的表」——更常见的是几块表格
    上下堆着、中间空几行。把整个 sheet 当一张表读，空行会把列对齐冲掉，
    两块不相干的数据也会被切进同一个 chunk。
    """
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "年度汇总"

    # 区域一：A1:C5
    ws["A1"] = "一季度收入"
    ws.merge_cells("A1:C1")
    for col, text in zip("ABC", ("月份", "收入", "成本"), strict=True):
        ws[f"{col}2"] = text
    for i, (month, rev, cost) in enumerate(
        [("1月", 320000, 210000), ("2月", 280000, 195000), ("3月", 410000, 260000)], start=3
    ):
        ws[f"A{i}"], ws[f"B{i}"], ws[f"C{i}"] = month, rev, cost

    # 空三行后，区域二：A9:D14（列数与区域一不同，压平就会串列）
    ws["A9"] = "重点客户明细"
    ws.merge_cells("A9:D9")
    for col, text in zip("ABCD", ("客户", "行业", "合同额", "签约日"), strict=True):
        ws[f"{col}10"] = text
    clients = [
        ("华东能源集团", "能源", 1850000, "2026-02-11"),
        ("南方物流", "物流", 920000, "2026-03-05"),
        ("中原制造", "制造", 1340000, "2026-03-28"),
    ]
    for i, row in enumerate(clients, start=11):
        for col, val in zip("ABCD", row, strict=True):
            ws[f"{col}{i}"] = val

    # 区域三：H2:I5，与前两块横向错开（同行不同列的独立区块）
    ws["H2"] = "指标"
    ws["I2"] = "数值"
    for i, (name, val) in enumerate(
        [("毛利率", "34.2%"), ("回款率", "91.5%"), ("客户数", 47)], start=3
    ):
        ws[f"H{i}"], ws[f"I{i}"] = name, val

    wb.save(path)


# ---------------------------------------------------------------- pdf


def make_text_pdf(path: Path) -> None:
    """数字 PDF（有文本层）：验证文本覆盖率与标题启发式。"""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    font_name = _register_cjk_font()
    cn = font_name == "CJK"
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4

    def page(title: str, body: list[str]) -> None:
        c.setFont(font_name, 18)
        c.drawString(60, height - 80, title)
        c.setFont(font_name, 11)
        y = height - 120
        for line in body:
            c.drawString(60, y, line)
            y -= 18
        c.showPage()

    if cn:
        page("技术方案说明书", [
            "本方案面向内网环境下的文档知识库建设。",
            "系统完全离线运行，不依赖任何外部服务。",
            "解析引擎支持七种常见办公文档格式。",
        ])
        page("第二章 实施计划", [
            "第一阶段完成环境准备与数据迁移。",
            "第二阶段完成解析入库与质量验收。",
            "交付期限为合同签订后 90 天。",
        ])
    else:
        page("Technical Proposal", [
            "This proposal targets offline knowledge base construction.",
            "The system runs fully offline with no external services.",
            "The parsing engine supports seven office document formats.",
        ])
        page("Chapter 2 Implementation", [
            "Phase one covers environment setup and data migration.",
            "Phase two covers parsing, ingestion and acceptance.",
            "Delivery deadline is 90 days after contract signing.",
        ])
    c.save()


def make_two_column_pdf(path: Path) -> None:
    """双栏排版 PDF：08 章 §3.2 的「多栏阅读顺序正确率 ≥ 90%」指标载体。

    多栏是本产品卖点最直接的试金石：按坐标从上到下扫，左右栏会交替串成
    「左1 右1 左2 右2」的乱序文本，读起来完全不成句。这份样本的每一栏
    自成完整段落，标注只要断言两栏各自的句子完整出现即可 —— 顺序正确率
    等 M2 接入版面分析后再收紧为逐句序号断言。
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate, Paragraph

    font_name = _register_cjk_font()
    cn = font_name == "CJK"
    width, height = A4
    margin, gutter = 50, 20
    col_w = (width - 2 * margin - gutter) / 2

    doc = BaseDocTemplate(str(path), pagesize=A4, leftMargin=margin, rightMargin=margin)
    frames = [
        Frame(margin, margin, col_w, height - 2 * margin, id="left"),
        Frame(margin + col_w + gutter, margin, col_w, height - 2 * margin, id="right"),
    ]
    doc.addPageTemplates([PageTemplate(id="two-col", frames=frames)])

    body = ParagraphStyle("body", fontName=font_name, fontSize=10, leading=15)
    title = ParagraphStyle("title", fontName=font_name, fontSize=15, leading=20, spaceAfter=10)

    if cn:
        blocks = [
            ("摘要", "本文提出一种面向内网环境的离线文档解析方法，"
                     "在不依赖任何外部服务的前提下完成结构化抽取。"),
            ("方法", "解析流程分为三级降级链，逐页选择最合适的解析级别，"
                     "并对每一页记录所用级别与置信度。"),
            ("实验", "在包含五十份复杂文档的测试集上，本方法的文本覆盖率"
                     "达到百分之九十七以上，表格结构还原优于基线方法。"),
            ("结论", "实验表明该方法在完全离线的约束下仍能保持内容完整性，"
                     "适用于涉密与内网场景的知识库建设。"),
        ]
    else:
        blocks = [
            ("Abstract", "We present an offline document parsing method for intranet "
                         "environments that performs structured extraction without external services."),
            ("Method", "The pipeline uses a three-level degradation chain, selecting the "
                       "most suitable level per page and recording the level and confidence."),
            ("Experiments", "On a test set of fifty complex documents, our method reaches "
                            "a text coverage above ninety-seven percent, outperforming the baseline."),
            ("Conclusion", "Results show the method preserves content integrity under a "
                           "fully offline constraint, suitable for classified intranet deployments."),
        ]

    story = []
    for heading, text in blocks:
        story.append(Paragraph(heading, title))
        # 重复正文把每栏撑满，逼出真正的跨栏排布（一屏放得下就测不到多栏）
        for _ in range(3):
            story.append(Paragraph(text, body))
    doc.build(story)


def make_long_table_pdf(path: Path) -> None:
    """跨页长表 + 表头重复：验证跨页表合并（FR-02）与切片时的表头复制（04 章 §3.1）。

    跨页表被切断后，第二页那截会变成一堆没有表头的裸数字 —— 对 RAG 来说
    等于噪声。repeatRows=1 让 reportlab 在每页重画表头，这正是解析器需要
    识别并**合并**的形态（否则会当成两张独立的表）。
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import LongTable, Paragraph, SimpleDocTemplate, TableStyle

    font_name = _register_cjk_font()
    cn = font_name == "CJK"

    doc = SimpleDocTemplate(str(path), pagesize=A4, topMargin=50, bottomMargin=50)
    title = ParagraphStyle("t", fontName=font_name, fontSize=15, leading=20, spaceAfter=12)

    header = (["序号", "物料编码", "名称", "数量", "单价"] if cn
              else ["No.", "Code", "Name", "Qty", "Price"])
    name = "配件" if cn else "Part"
    rows = [[str(i), f"WL-{i:05d}", f"{name}{i:03d}", str(i % 20 + 1), f"{80 + i * 1.5:.2f}"]
            for i in range(1, 121)]

    table = LongTable([header, *rows], repeatRows=1, colWidths=[50, 90, 130, 50, 70])
    table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, -1), font_name),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
        ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
    ]))
    doc.build([Paragraph("物料清单（跨页表）" if cn else "Bill of Materials", title), table])


def make_scanned_pdf(path: Path) -> None:
    """无文本层的图片 PDF：模拟扫描件，走 L2 整页 OCR 分支。

    **当前实现下这份样本的文本覆盖率必然是 0** —— RapidOCR 尚未接入（03 章 §3，
    M2 交付）。标注因此只断言「不崩溃、有明确产出」，不设覆盖率下限。
    M2 接入 OCR 后必须回来把 min_text_coverage 收紧到 0.90，
    否则这份样本就只是在证明「引擎没炸」，测不到 OCR 到底有没有生效。
    """
    from io import BytesIO

    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    width, height = A4
    c = canvas.Canvas(str(path), pagesize=A4)
    for page, rgb in enumerate([(238, 238, 230), (232, 236, 240)], start=1):
        # 整页铺一张纯色位图：PDF 里没有任何字形，pdfplumber/pypdfium2 提不出一个字符
        c.drawImage(ImageReader(BytesIO(_png_bytes(300, 420, rgb))),
                    0, 0, width=width, height=height)
        # 画几条深色横杠模拟文字行，让它在预览里也像一页扫描件
        c.setFillColorRGB(0.25, 0.25, 0.28)
        for i in range(12):
            y = height - 120 - i * 26
            c.rect(70, y, (width - 200) * (0.55 + 0.4 * ((i * 7 + page) % 5) / 5), 9, fill=1, stroke=0)
        c.showPage()
    c.save()


# ---------------------------------------------------------------- 异常样本


def make_encrypted_pdf(path: Path) -> None:
    """带打开口令的 PDF：验证 E02「文件受密码保护」的优雅报错（100% 门禁）。

    01 章 §2.3 明确把「密码破解」列进不做清单，所以正确行为是**干净地报 E02**
    并告诉用户去掉密码后重试，而不是尝试空口令、更不是抛栈。
    """
    from reportlab.lib import pdfencrypt
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    font_name = _register_cjk_font()
    enc = pdfencrypt.StandardEncryption("userpass2026", ownerPassword="ownerpass2026")
    c = canvas.Canvas(str(path), pagesize=A4, encrypt=enc)
    c.setFont(font_name, 14)
    c.drawString(60, A4[1] - 90, "机密文件" if font_name == "CJK" else "Confidential")
    c.setFont(font_name, 11)
    c.drawString(60, A4[1] - 130,
                 "此文件受口令保护。" if font_name == "CJK" else "This file is password protected.")
    c.showPage()
    c.save()


def make_corrupt_docx(path: Path) -> None:
    """损坏文件：字节是 zip 但内部不是 OOXML —— 必须报 E01 而不是崩溃。"""
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("not-a-document.txt", "这不是一个有效的 docx")


def make_truncated_pdf(path: Path) -> None:
    """截断的 PDF：只有文件头，验证 E01 优雅报错。"""
    path.write_bytes(b"%PDF-1.7\n1 0 obj\n<< /Type /Catalog")


def make_unsupported_file(path: Path) -> None:
    """非支持格式：导入阶段就该被 probe_file 挡下（E03）。"""
    path.write_text("plain text, not an office document", encoding="utf-8")


# ---------------------------------------------------------------- 注册表

FIXTURES: dict[str, Callable[[Path], None]] = {
    # —— 结构还原 ——
    "headings.docx": make_headings_docx,
    "merged_table.docx": make_merged_table_docx,
    "nested_list.docx": make_nested_list_docx,
    "header_footer.docx": make_header_footer_docx,
    "with_image.docx": make_image_docx,
    "bilingual.docx": make_bilingual_docx,
    "slides.pptx": make_slides_pptx,
    "multisheet.xlsx": make_multisheet_xlsx,
    "sparse_regions.xlsx": make_sparse_regions_xlsx,
    "legacy_table.xls": make_legacy_table_xls,
    "two_column.pdf": make_two_column_pdf,
    "long_table.pdf": make_long_table_pdf,
    "text_document.pdf": make_text_pdf,
    # —— 规模与退化路径 ——
    "wide_table.xlsx": make_wide_xlsx,
    "large_deck.pptx": make_large_deck_pptx,
    "scanned.pdf": make_scanned_pdf,
    # —— 优雅报错（有产出不是唯一的「通过」）——
    "empty.docx": make_empty_docx,
    "corrupt.docx": make_corrupt_docx,
    "truncated.pdf": make_truncated_pdf,
    "encrypted.pdf": make_encrypted_pdf,
    "unsupported.txt": make_unsupported_file,
}


def generate(out_dir: Path = DEFAULT_OUT, *, force: bool = False) -> list[Path]:
    """生成全部 fixtures，返回产出路径。已存在且非 force 时跳过（真实语料不会被覆盖）。"""
    out_dir.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    for name, fn in FIXTURES.items():
        target = out_dir / name
        if target.exists() and not force:
            continue
        try:
            fn(target)
            written.append(target)
        except ImportError as exc:
            print(f"跳过 {name}：缺少依赖（{exc}）", file=sys.stderr)
        except Exception as exc:
            print(f"生成 {name} 失败：{type(exc).__name__}: {exc}", file=sys.stderr)
    return written


def main(argv: list[str] | None = None) -> int:
    args = argv if argv is not None else sys.argv[1:]
    force = "--force" in args
    positional = [a for a in args if not a.startswith("--")]
    out_dir = Path(positional[0]) if positional else DEFAULT_OUT

    written = generate(out_dir, force=force)
    print(f"已生成 {len(written)} 个样本到 {out_dir}")
    for p in written:
        print(f"  + {p.name}")
    if not written:
        print("（样本已存在，加 --force 可重新生成）")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
